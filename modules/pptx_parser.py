"""
FIDELITAS — PPTX Reference Parser

A client PPTX design deck has no ▼ALT/▼CTA marker convention, but it does
have real, extractable structure once you look at shape types rather than
just raw text:

  - A metadata/brief slide using "■Label：value" or "◼︎Label：value" lines
    (label and value sometimes on the same line, sometimes value on the
    next line) — "■メール件名：" is the mail subject line specifically,
    extracted as kind="title" so it flows through the same
    content_engine.check_mail_title() used for the checklist workbook.
  - CTA button mockups are AUTO_SHAPE shapes with short, single-line text
    (e.g. "詳細はこちら", "BMW BELIEVES") — extracted as kind="pptx_cta_text".
  - Bare-URL text boxes near those buttons — extracted as kind="pptx_url".
    IMPORTANT: these are NOT paired 1:1 to a specific button. Checked
    against real slide data, shape order is NOT reliable for that pairing
    (on one slide the URL follows immediately; on another, several image
    shapes sit between a button and its URL, because slides get edited
    non-linearly). Claiming a specific button->URL pairing the deck itself
    doesn't reliably encode would be fabricating a result, so button text
    and URLs are each checked as independently-verifiable presence facts
    instead: "does this button copy appear in the HTML" and "does this URL
    appear in the HTML" — both true statements, without inventing a link
    between them that isn't actually there.
  - Boilerplate (the same exact text repeated on 3+ slides — page
    watermarks like "DEM構成案｜..." ) is detected and excluded, since it's
    deck chrome, not campaign content, and would otherwise generate a
    guaranteed-false "not found in HTML" result on every single slide.
  - Everything else with a text frame is body copy, presence-checked.
  - Speaker notes, and a per-slide image count, are captured as before.
"""

import re
from collections import Counter
from pptx import Presentation
from .models import ReferenceModel, ReferenceVariant, ReferenceBlock

_FIELD_LINE_RE = re.compile(r"^[■◼︎\ufe0e\ufe0f]+\s*([^：:]+)[：:]\s*(.*)$")
_URL_RE = re.compile(r"^https?://\S+$")
_BOILERPLATE_MIN_REPEATS = 3
_CTA_MAX_LEN = 40


def parse_pptx(path: str) -> tuple[ReferenceModel, list]:
    warnings = []
    model = ReferenceModel()
    try:
        prs = Presentation(path)
    except Exception as e:
        warnings.append(f"Could not open this file as a PPTX: {e}")
        return model, warnings

    variant = ReferenceVariant(name="PPTX Reference")
    order = 0

    # ---- pass 1: collect every text frame, to detect cross-slide boilerplate ----
    collected = []  # (slide_idx, text, shape_type)
    image_counts = {}
    for slide_idx, slide in enumerate(prs.slides, start=1):
        count = 0
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                collected.append((slide_idx, shape.text_frame.text.strip(), str(shape.shape_type)))
            if "PICTURE" in str(shape.shape_type):
                count += 1
        image_counts[slide_idx] = count

    text_counts = Counter(t[1] for t in collected)
    boilerplate = {t for t, c in text_counts.items() if c >= _BOILERPLATE_MIN_REPEATS}

    # ---- pass 2: classify each non-boilerplate text frame ----
    for slide_idx, text, shape_type in collected:
        if text in boilerplate:
            continue

        lines = text.split("\n")
        field_hits = sum(1 for l in lines if _FIELD_LINE_RE.match(l))

        if field_hits >= 2:
            # metadata/brief block — parse field-by-field rather than as body text
            i = 0
            while i < len(lines):
                m = _FIELD_LINE_RE.match(lines[i])
                if m:
                    label, inline_val = m.group(1).strip(), m.group(2).strip()
                    val = inline_val
                    if not val:
                        j = i + 1
                        while j < len(lines) and not lines[j].strip():
                            j += 1
                        if j < len(lines) and not _FIELD_LINE_RE.match(lines[j]):
                            val = lines[j].strip()
                            i = j
                    if val:
                        if "件名" in label:
                            order += 1
                            variant.blocks.append(ReferenceBlock(order=order, kind="title", text=val))
                        elif "s3" in label.lower() or "link" in label.lower():
                            order += 1
                            variant.blocks.append(ReferenceBlock(
                                order=order, kind="pptx_reference_url",
                                text=f"{label}: {val}", url=val,
                            ))
                        else:
                            order += 1
                            variant.blocks.append(ReferenceBlock(
                                order=order, kind="pptx_notes", text=f"{label}: {val}",
                            ))
                i += 1
            continue

        if shape_type.startswith("AUTO_SHAPE") and len(text) <= _CTA_MAX_LEN and "\n" not in text:
            order += 1
            variant.blocks.append(ReferenceBlock(order=order, kind="pptx_cta_text", text=text))
            continue

        if _URL_RE.match(text):
            order += 1
            variant.blocks.append(ReferenceBlock(order=order, kind="pptx_url", text=text, url=text))
            continue

        order += 1
        variant.blocks.append(ReferenceBlock(order=order, kind="pptx_body", text=text))

    # ---- speaker notes + per-slide image counts ----
    for slide_idx, slide in enumerate(prs.slides, start=1):
        try:
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    order += 1
                    variant.blocks.append(ReferenceBlock(order=order, kind="pptx_notes", text=notes_text))
        except Exception:
            pass
        if image_counts.get(slide_idx):
            order += 1
            variant.blocks.append(ReferenceBlock(
                order=order, kind="pptx_image_count",
                text=f"Slide {slide_idx}: {image_counts[slide_idx]} image(s)",
            ))

    if not variant.blocks:
        warnings.append("No text or notes could be extracted from any slide.")
    if boilerplate:
        warnings.append(f"Excluded {len(boilerplate)} repeated boilerplate text block(s) "
                         f"(appeared on {_BOILERPLATE_MIN_REPEATS}+ slides — treated as deck chrome, not content).")

    model.variants.append(variant)
    return model, warnings
